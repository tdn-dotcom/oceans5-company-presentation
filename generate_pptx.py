"""
Oceans5 Company Presentation — PPTX Generator
Generates Oceans5_Company_Presentation_2026_02_English.pptx
Run: python generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import os

# ── Brand colors ──────────────────────────────────────────────
NAVY_DEEP    = RGBColor(0x0A, 0x1B, 0x2E)
NAVY_MID     = RGBColor(0x12, 0x28, 0x48)
BLUE_PRIMARY = RGBColor(0x1A, 0x52, 0x76)
TEAL         = RGBColor(0x00, 0x7E, 0xA7)
TEAL_BRIGHT  = RGBColor(0x00, 0xB4, 0xD8)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE    = RGBColor(0xF4, 0xF8, 0xFB)
GREY_LIGHT   = RGBColor(0xB8, 0xC8, 0xD8)
GREY_MID     = RGBColor(0x6B, 0x8B, 0xA4)
TEXT_DARK    = RGBColor(0x0A, 0x1B, 0x2E)
GOLD         = RGBColor(0xF0, 0xA5, 0x00)

# ── Slide dimensions: 16:9 widescreen ─────────────────────────
W = Inches(13.333)
H = Inches(7.5)

OUT = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "Oceans5_Company_Presentation_2026_02_English.pptx"
)

# ── Helper: new presentation ───────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

# ── Helper: blank slide ────────────────────────────────────────
def add_blank(prs):
    layout = prs.slide_layouts[6]   # blank
    return prs.slides.add_slide(layout)

# ── Helper: solid fill a shape ────────────────────────────────
def fill_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb

# ── Helper: add rectangle ─────────────────────────────────────
def add_rect(slide, l, t, w, h, rgb):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    fill_solid(shape, rgb)
    shape.line.fill.background()
    return shape

# ── Helper: add text box ───────────────────────────────────────
def add_tb(slide, text, l, t, w, h,
           size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
           font_name="Montserrat"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = font_name
    return txb

# ── Helper: accent line (teal bar) ────────────────────────────
def accent_bar(slide, l, t):
    bar = add_rect(slide, l, t, Inches(0.7), Pt(3), TEAL_BRIGHT)
    return bar

# ── Helper: tag label ─────────────────────────────────────────
def add_tag(slide, text, l, t):
    add_tb(slide, text.upper(), l, t, Inches(6), Inches(0.3),
           size=9, bold=True, color=TEAL_BRIGHT, font_name="Inter")

# ── Helper: footer ─────────────────────────────────────────────
def add_footer(slide, left_text="", right_text="© 2026 Oceans5 Group AS", dark=False):
    col = RGBColor(0x44, 0x55, 0x66) if dark else RGBColor(0x88, 0x99, 0xAA)
    add_tb(slide, left_text, Inches(0.6), Inches(7.15), Inches(8), Inches(0.3),
           size=8, color=col, font_name="Inter")
    add_tb(slide, right_text, Inches(5.8), Inches(7.15), Inches(7.3), Inches(0.3),
           size=8, color=col, align=PP_ALIGN.RIGHT, font_name="Inter")

# ── Helper: card background ────────────────────────────────────
def add_card(slide, l, t, w, h, dark=True):
    rgb = RGBColor(0x16, 0x2E, 0x4A) if dark else RGBColor(0xFF, 0xFF, 0xFF)
    r = add_rect(slide, l, t, w, h, rgb)
    return r

# ═══════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    s = add_blank(prs)
    # Background gradient approximated with two rects
    add_rect(s, 0, 0, W, H, NAVY_DEEP)
    add_rect(s, Inches(7), 0, Inches(6.333), H, RGBColor(0x00, 0x5A, 0x75))

    # Logo text
    add_tb(s, "Oceans5", Inches(4.2), Inches(1.4), Inches(5), Inches(1.2),
           size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Montserrat")
    add_rect(s, Inches(4.2), Inches(2.55), Inches(4.6), Pt(3), TEAL_BRIGHT)

    # Tagline
    add_tb(s, "Engineering Tomorrow's Ocean Economy",
           Inches(1.5), Inches(2.9), Inches(10.3), Inches(1.2),
           size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Montserrat")

    # Sub
    add_tb(s, "Marine Engineering  ·  AI Project Management  ·  Offshore Consultancy",
           Inches(2), Inches(4.0), Inches(9.3), Inches(0.6),
           size=13, color=GREY_LIGHT, align=PP_ALIGN.CENTER, font_name="Inter")

    # Stats row
    stats = [("2021", "Founded"), ("Oslo", "Norway"), ("3", "Companies"), ("5", "Senior Owners")]
    for i, (val, lbl) in enumerate(stats):
        x = Inches(2.5 + i * 2.2)
        add_tb(s, val, x, Inches(5.0), Inches(2), Inches(0.7),
               size=26, bold=True, color=TEAL_BRIGHT, align=PP_ALIGN.CENTER, font_name="Montserrat")
        add_tb(s, lbl, x, Inches(5.65), Inches(2), Inches(0.4),
               size=10, color=GREY_LIGHT, align=PP_ALIGN.CENTER, font_name="Inter")

    add_footer(s, "© 2026 Oceans5 Group AS · oceans5.no", "Confidential")
    return s


def slide_02_challenge(prs):
    s = add_blank(prs)
    add_rect(s, 0, 0, W, H, NAVY_DEEP)

    add_tag(s, "The Challenge", Inches(0.6), Inches(0.4))
    add_tb(s, "Offshore projects are getting harder, faster — and riskier",
           Inches(0.6), Inches(0.7), Inches(12), Inches(1.0),
           size=30, bold=True, color=WHITE, font_name="Montserrat")
    accent_bar(s, Inches(0.6), Inches(1.65))

    cards = [
        ("⚡ Complexity at scale",
         "Offshore wind, HVDC cables, and marine installations now require\n"
         "cross-discipline expertise spanning engineering, regulation,\n"
         "logistics, and AI — simultaneously. Few teams can do it all."),
        ("💸 Cost & schedule overruns",
         "Major offshore projects routinely run 20–40% over budget.\n"
         "Fragmented supply chains, poor data flow, and inadequate\n"
         "engineering oversight cost billions across the industry each year."),
        ("🤖 AI is reshaping the rules",
         "Clients are adopting AI faster than traditional consultancies can\n"
         "respond. Companies that don't embed AI into every workflow risk\n"
         "becoming obsolete — within years, not decades."),
    ]
    for i, (title, body) in enumerate(cards):
        x = Inches(0.4 + i * 4.3)
        add_card(s, x, Inches(2.0), Inches(4.1), Inches(4.0))
        add_tb(s, title, x + Inches(0.2), Inches(2.1), Inches(3.7), Inches(0.6),
               size=14, bold=True, color=TEAL_BRIGHT, font_name="Montserrat")
        add_tb(s, body, x + Inches(0.2), Inches(2.7), Inches(3.7), Inches(3.0),
               size=11, color=GREY_LIGHT, font_name="Inter")

    # Bottom callout
    add_rect(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.85), RGBColor(0x00, 0x48, 0x62))
    add_tb(s, "The market gap: Project owners need one trusted partner who combines deep marine engineering\n"
              "expertise with AI-augmented delivery — not three separate vendors and a coordination headache.",
           Inches(0.6), Inches(6.25), Inches(12.1), Inches(0.75),
           size=11, color=WHITE, font_name="Inter")

    add_footer(s, "Oceans5 · The Challenge")
    return s


def slide_03_who_we_are(prs):
    s = add_blank(prs)
    add_rect(s, 0, 0, W, H, NAVY_DEEP)

    add_tag(s, "Who We Are", Inches(0.6), Inches(0.4))
    add_tb(s, "One organisation. Concept to completion.",
           Inches(0.6), Inches(0.7), Inches(12), Inches(0.9),
           size=30, bold=True, color=WHITE, font_name="Montserrat")
    accent_bar(s, Inches(0.6), Inches(1.55))

    body = (
        "Oceans5 is a Norwegian specialist firm founded in 2021. We combine engineering services,\n"
        "AI-augmented project management, and operational execution in one organisation — so clients\n"
        "meet the same team from concept to completed installation.\n\n"
        "We work primarily in marine installations, offshore wind, power cable infrastructure,\n"
        "and advisory services — with growing focus on AI-driven delivery and defined work packages."
    )
    add_tb(s, body, Inches(0.6), Inches(1.75), Inches(8.2), Inches(2.4),
           size=12, color=GREY_LIGHT, font_name="Inter")

    sub_cards = [
        ("Independent", "No vendor bias — we advise on what is technically right."),
        ("AI-First", "AI embedded in every delivery and internal workflow."),
        ("Agile & premium", "Small enough to move fast. Experienced enough to get it right."),
        ("End-to-end", "Engineering, project management, and crew — one team."),
    ]
    for i, (t, b) in enumerate(sub_cards):
        col = i % 2
        row = i // 2
        x = Inches(0.4 + col * 4.2)
        y = Inches(4.0 + row * 1.3)
        add_card(s, x, y, Inches(4.0), Inches(1.15))
        add_tb(s, t, x + Inches(0.15), y + Inches(0.1), Inches(3.7), Inches(0.45),
               size=13, bold=True, color=TEAL_BRIGHT, font_name="Montserrat")
        add_tb(s, b, x + Inches(0.15), y + Inches(0.52), Inches(3.7), Inches(0.55),
               size=10, color=GREY_LIGHT, font_name="Inter")

    # Right metrics
    metrics = [("2021", "Founded in Oslo, Norway"), ("3", "Companies in the group"),
               ("5", "Senior co-founder owners"), ("Nordic+", "Core market")]
    for i, (val, lbl) in enumerate(metrics):
        y = Inches(1.75 + i * 1.35)
        add_card(s, Inches(9.1), y, Inches(3.8), Inches(1.2))
        add_tb(s, val, Inches(9.1), y + Inches(0.1), Inches(3.8), Inches(0.7),
               size=28, bold=True, color=TEAL_BRIGHT, align=PP_ALIGN.CENTER, font_name="Montserrat")
        add_tb(s, lbl, Inches(9.1), y + Inches(0.8), Inches(3.8), Inches(0.35),
               size=9, color=GREY_LIGHT, align=PP_ALIGN.CENTER, font_name="Inter")

    add_footer(s, "Oceans5 · Who We Are")
    return s


def slide_04_services(prs):
    s = add_blank(prs)
    add_rect(s, 0, 0, W, H, OFF_WHITE)

    add_tag(s, "What We Do", Inches(0.6), Inches(0.4))
    add_tb(s, "Four integrated service lines",
           Inches(0.6), Inches(0.7), Inches(12), Inches(0.9),
           size=30, bold=True, color=TEXT_DARK, font_name="Montserrat")
    accent_bar(s, Inches(0.6), Inches(1.55))

    services = [
        ("⚓", "Marine Operations",
         "Project planning and execution for complex marine operations. "
         "Installation engineering, load-out and transport, route planning, "
         "cable protection, and contingency management."),
        ("🔌", "Power & Cable Engineering",
         "Full-scope engineering for offshore wind, export cables, and "
         "interconnectors. From pre-FID concept to construction support "
         "and as-built documentation."),
        ("🗺️", "Survey & Route Engineering",
         "Integrated survey and engineering from pre-route through as-built. "
         "Pre-route, boulder, UXO, seabed mobility, as-laid, intermediate, "
         "and O&M surveys with full engineering analysis."),
        ("🧠", "AI Project Management",
         "AI-augmented delivery across all project phases. From AI-assisted "
         "reporting and risk analysis to intelligent schedule management and "
         "real-time operational monitoring."),
    ]
    for i, (icon, title, body) in enumerate(services):
        x = Inches(0.3 + i * 3.2)
        add_rect(s, x, Inches(1.9), Inches(3.05), Inches(4.6), WHITE)
        add_tb(s, icon, x + Inches(0.2), Inches(2.0), Inches(0.6), Inches(0.6), size=20, color=TEAL)
        add_tb(s, title, x + Inches(0.2), Inches(2.55), Inches(2.7), Inches(0.55),
               size=13, bold=True, color=TEAL, font_name="Montserrat")
        add_tb(s, body, x + Inches(0.2), Inches(3.1), Inches(2.75), Inches(3.0),
               size=10, color=TEXT_DARK, font_name="Inter")

    add_rect(s, Inches(0.3), Inches(6.65), Inches(12.7), Inches(0.6), RGBColor(0xE8, 0xF4, 0xF8))
    add_tb(s, "Work packages, not just people: We are moving beyond hourly resource supply to structured, "
              "outcome-based packages — FEED studies, survey deliveries, and Client rep services with defined scope, "
              "price, and accountability.",
           Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.55),
           size=10, color=TEAL, font_name="Inter")

    add_footer(s, "Oceans5 · Services", dark=True)
    return s


def slide_05_ai_advantage(prs):
    s = add_blank(prs)
    add_rect(s, 0, 0, W, H, NAVY_MID)

    add_tag(s, "The AI Advantage", Inches(0.6), Inches(0.4))
    add_tb(s, "AI isn't a tool we use. It's how we work.",
           Inches(0.6), Inches(0.7), Inches(12), Inches(0.9),
           size=30, bold=True, color=WHITE, font_name="Montserrat")
    accent_bar(s, Inches(0.6), Inches(1.55))

    pillars = [
        ("01", "Every workflow, AI-enabled",
         "From proposal writing and engineering reports to schedule management, "
         "risk analysis, and client communication — AI compresses timelines and eliminates rework."),
        ("02", "Distributed ownership of innovation",
         "Every Oceans5 engineer owns an AI experiment. New tools, new workflows — proposed "
         "from the floor, piloted in 90 days, scaled or killed based on results."),
        ("03", "AI-augmented project management",
         "We combine senior engineering judgment with AI monitoring to catch schedule risk early, "
         "surface change order patterns, and keep clients informed in real time."),
        ("04", "Applied — not hypothetical",
         "We do not sell AI strategy decks. We use AI daily, measure it, and bring that live "
         "experience to every client engagement. This is our competitive moat."),
    ]
    for i, (num, title, body) in enumerate(pillars):
        y = Inches(1.85 + i * 1.25)
        add_tb(s, num, Inches(0.4), y, Inches(0.7), Inches(0.8),
               size=28, bold=True, color=TEAL_BRIGHT, font_name="Montserrat")
        add_tb(s, title, Inches(1.1), y, Inches(5.8), Inches(0.45),
               size=13, bold=True, color=WHITE, font_name="Montserrat")
        add_tb(s, body, Inches(1.1), y + Inches(0.42), Inches(5.8), Inches(0.75),
               size=10, color=GREY_LIGHT, font_name="Inter")

    # Right side: AI chart (using pptx bar chart)
    chart_data = ChartData()
    chart_data.categories = ['Report gen.', 'Risk analysis', 'Scheduling', 'Change orders', 'Client comms']
    chart_data.add_series('Traditional', (100, 100, 100, 100, 100))
    chart_data.add_series('AI-augmented (Oceans5)', (45, 60, 55, 65, 50))

    chart = s.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(7.2), Inches(1.8), Inches(5.8), Inches(4.2),
        chart_data
    )
    c = chart.chart
    c.has_legend = True
    c.legend.position = 2  # bottom

    try:
        ser0 = c.series[0]
        ser0.format.fill.solid()
        ser0.format.fill.fore_color.rgb = GREY_MID
        ser1 = c.series[1]
        ser1.format.fill.solid()
        ser1.format.fill.fore_color.rgb = TEAL_BRIGHT
    except Exception:
        pass

    add_rect(s, Inches(7.0), Inches(6.2), Inches(6.0), Inches(0.9), RGBColor(0x00, 0x35, 0x50))
    add_tb(s, '"By 2028 — we are recognised as the most AI-driven engineering partner in the Nordic offshore market."',
           Inches(7.1), Inches(6.25), Inches(5.9), Inches(0.8),
           size=10, color=WHITE, font_name="Inter")

    add_footer(s, "Oceans5 · The AI Advantage")
    return s


def slide_06_track_record(prs):
    s = add_blank(prs)
    add_rect(s, 0, 0, W, H, NAVY_DEEP)

    add_tag(s, "Track Record", Inches(0.6), Inches(0.4))
    add_tb(s, "Proven on the world's most demanding offshore projects",
           Inches(0.6), Inches(0.7), Inches(12.2), Inches(0.9),
           size=30, bold=True, color=WHITE, font_name="Montserrat")
    accent_bar(s, Inches(0.6), Inches(1.55))

    # Cable projects
    add_card(s, Inches(0.4), Inches(1.8), Inches(6.1), Inches(2.4))
    add_tb(s, "⚡ Cable & Power Projects", Inches(0.6), Inches(1.9), Inches(5.8), Inches(0.5),
           size=13, bold=True, color=TEAL_BRIGHT, font_name="Montserrat")
    projects_cable = ("Empire Wind  ·  Estlink  ·  Revolution  ·  Sunrise\n"
                      "Tyrrhenian Link  ·  Borwin5  ·  Dogger Bank\n"
                      "East Anglia  ·  Baltyk 2 & 3  ·  Beatrice (repair)  ·  Elia (repair)")
    add_tb(s, projects_cable, Inches(0.6), Inches(2.4), Inches(5.7), Inches(1.7),
           size=11, color=GREY_LIGHT, font_name="Inter")

    # Marine ops
    add_card(s, Inches(0.4), Inches(4.3), Inches(6.1), Inches(2.1))
    add_tb(s, "⚓ Marine Operations", Inches(0.6), Inches(4.4), Inches(5.8), Inches(0.5),
           size=13, bold=True, color=TEAL_BRIGHT, font_name="Montserrat")
    projects_marine = ("Wilhelmshaven LNG Terminal  ·  Fensfjorden\n"
                       "New-build CLV Inspection (Turkey)\n"
                       "Cable-laying Equipment Design  ·  Marine Mooring Studies")
    add_tb(s, projects_marine, Inches(0.6), Inches(4.9), Inches(5.7), Inches(1.4),
           size=11, color=GREY_LIGHT, font_name="Inter")

    # Advisory
    add_card(s, Inches(6.8), Inches(1.8), Inches(6.1), Inches(2.4))
    add_tb(s, "📐 Advisory & Engineering", Inches(7.0), Inches(1.9), Inches(5.8), Inches(0.5),
           size=13, bold=True, color=TEAL_BRIGHT, font_name="Montserrat")
    add_tb(s, ("Pre-FID Concept Studies\nThird-Party Reviews\n"
               "Constructability Reviews\nTechnical Incident Advisory"),
           Inches(7.0), Inches(2.4), Inches(5.7), Inches(1.7),
           size=11, color=GREY_LIGHT, font_name="Inter")

    # Geography
    add_card(s, Inches(6.8), Inches(4.3), Inches(6.1), Inches(2.1))
    add_tb(s, "🌍 Geographic Coverage", Inches(7.0), Inches(4.4), Inches(5.8), Inches(0.5),
           size=13, bold=True, color=TEAL_BRIGHT, font_name="Montserrat")
    add_tb(s, ("Projects spanning the North Sea, Baltic, Mediterranean,\n"
               "and US Atlantic — primary focus on Nordic and\n"
               "Western European offshore market."),
           Inches(7.0), Inches(4.9), Inches(5.7), Inches(1.4),
           size=11, color=GREY_LIGHT, font_name="Inter")

    add_footer(s, "Oceans5 · Track Record")
    return s


def slide_07_approach(prs):
    s = add_blank(prs)
    add_rect(s, 0, 0, W, H, OFF_WHITE)

    add_tag(s, "How We Work", Inches(0.6), Inches(0.4))
    add_tb(s, "Agile by design. Methodical by discipline.",
           Inches(0.6), Inches(0.7), Inches(12), Inches(0.9),
           size=30, bold=True, color=TEXT_DARK, font_name="Montserrat")
    accent_bar(s, Inches(0.6), Inches(1.55))

    intro = ("We run in 90-day delivery cycles — structured, measurable, and deliberately iterative. "
             "Every pilot and work package is evaluated at cycle end: continue, adjust, or close.\n"
             "We reward learning, not self-preservation.")
    add_tb(s, intro, Inches(0.6), Inches(1.75), Inches(6.5), Inches(1.2),
           size=12, color=TEXT_DARK, font_name="Inter")

    steps = [
        ("1", "Early engagement & FEED", "Pre-FID concept, risk identification, scope definition"),
        ("2", "Engineering & survey", "Detailed design, route studies, regulatory compliance"),
        ("3", "AI-augmented execution", "Real-time monitoring, intelligent reporting, risk alerts"),
        ("4", "Client rep & handover", "Owner's representation, testing, commissioning, as-built"),
    ]
    for i, (n, t, b) in enumerate(steps):
        y = Inches(3.0 + i * 0.95)
        add_rect(s, Inches(0.4), y, Inches(0.45), Inches(0.45), TEAL)
        add_tb(s, n, Inches(0.4), y, Inches(0.45), Inches(0.45),
               size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Montserrat")
        add_tb(s, t, Inches(0.95), y + Inches(0.02), Inches(5.6), Inches(0.35),
               size=12, bold=True, color=TEXT_DARK, font_name="Montserrat")
        add_tb(s, b, Inches(0.95), y + Inches(0.38), Inches(5.6), Inches(0.35),
               size=10, color=GREY_MID, font_name="Inter")

    right_cards = [
        ("Work packages, not just people", "Defined scope · Defined price · Full accountability"),
        ("ISO 9001 pathway", "QMS in active use · Certification target 2026–2027"),
        ("90-day pilot cycle", "Run · Measure · Learn · Scale or close"),
    ]
    for i, (t, b) in enumerate(right_cards):
        y = Inches(1.8 + i * 1.7)
        add_rect(s, Inches(7.2), y, Inches(5.7), Inches(1.5), WHITE)
        add_tb(s, t, Inches(7.4), y + Inches(0.15), Inches(5.3), Inches(0.55),
               size=13, bold=True, color=TEAL, align=PP_ALIGN.CENTER, font_name="Montserrat")
        add_tb(s, b, Inches(7.4), y + Inches(0.7), Inches(5.3), Inches(0.55),
               size=10, color=GREY_MID, align=PP_ALIGN.CENTER, font_name="Inter")

    add_footer(s, "Oceans5 · How We Work", dark=True)
    return s


def slide_08_team(prs):
    s = add_blank(prs)
    add_rect(s, 0, 0, W, H, NAVY_DEEP)

    add_tag(s, "The Team", Inches(0.6), Inches(0.4))
    add_tb(s, "Senior experts. Co-owners with skin in the game.",
           Inches(0.6), Inches(0.7), Inches(12), Inches(0.9),
           size=30, bold=True, color=WHITE, font_name="Montserrat")
    accent_bar(s, Inches(0.6), Inches(1.55))

    team = [
        ("PEJ", "Peter-Emil Johannessen", "CEO & Co-founder",
         "Offshore electrification and wind power installation. Former CEO Havkonsult AS. Commercial leadership."),
        ("AMJ", "Anna M. E. Jarbekk", "CFO",
         "Project management and administration. Structures and improves operational processes."),
        ("MSM", "Morten S. Mathiesen", "QHSE Director & Co-founder",
         "Strategy, advisory and entrepreneurship. Drives ISO 9001 certification and QMS implementation."),
        ("TDN", "Thomas D. Nielsen", "Head of Engineering & Co-founder",
         "Offshore engineering and marine operations. AI strategy and digital enablement."),
        ("KA",  "Kent Aadnekvam", "Operations Lead",
         "Hands-on maritime and onshore operations. Experienced entrepreneur and skilled in execution."),
    ]
    for i, (initials, name, role, bio) in enumerate(team):
        x = Inches(0.25 + i * 2.6)
        add_card(s, x, Inches(1.9), Inches(2.5), Inches(4.6))
        # Avatar circle approximation
        add_rect(s, x + Inches(0.8), Inches(2.05), Inches(0.9), Inches(0.9), TEAL)
        add_tb(s, initials, x + Inches(0.8), Inches(2.05), Inches(0.9), Inches(0.9),
               size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Montserrat")
        add_tb(s, name, x + Inches(0.1), Inches(3.05), Inches(2.35), Inches(0.55),
               size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Montserrat")
        add_tb(s, role, x + Inches(0.1), Inches(3.55), Inches(2.35), Inches(0.45),
               size=9, color=TEAL_BRIGHT, align=PP_ALIGN.CENTER, font_name="Inter")
        add_tb(s, bio, x + Inches(0.1), Inches(3.95), Inches(2.35), Inches(2.2),
               size=9, color=GREY_LIGHT, align=PP_ALIGN.CENTER, font_name="Inter")

    add_rect(s, Inches(0.4), Inches(6.6), Inches(12.5), Inches(0.6), RGBColor(0x00, 0x35, 0x50))
    add_tb(s, "All five leaders are co-owners — our interests are fully aligned with yours. "
              "When your project succeeds, we succeed.",
           Inches(0.6), Inches(6.65), Inches(12.1), Inches(0.5),
           size=11, color=WHITE, align=PP_ALIGN.CENTER, font_name="Inter")

    add_footer(s, "Oceans5 · The Team")
    return s


def slide_09_market(prs):
    s = add_blank(prs)
    add_rect(s, 0, 0, W, H, OFF_WHITE)

    add_tag(s, "Market Opportunity", Inches(0.6), Inches(0.4))
    add_tb(s, "A once-in-a-generation energy buildout",
           Inches(0.6), Inches(0.7), Inches(12), Inches(0.9),
           size=30, bold=True, color=TEXT_DARK, font_name="Montserrat")
    accent_bar(s, Inches(0.6), Inches(1.55))

    # Chart
    chart_data = ChartData()
    chart_data.categories = ['2020', '2022', '2024', '2026E', '2028E', '2030E']
    chart_data.add_series('Global Offshore Wind (GW)', (35, 57, 90, 145, 230, 370))
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(0.4), Inches(1.8), Inches(6.8), Inches(4.5),
        chart_data
    )
    try:
        c = chart.chart
        c.has_legend = False
        ser = c.series[0]
        ser.format.line.color.rgb = TEAL
        ser.format.line.width = Pt(2.5)
    except Exception:
        pass

    market_cards = [
        ("€1 trillion+",
         "Planned offshore wind investment in Europe alone through 2030 — driven by REPowerEU and national targets."),
        ("~500 GW",
         "Offshore wind target for Europe by 2030. Every GW installed requires kilometres of cables and engineering oversight."),
        ("Skill gap = opportunity",
         "The offshore wind industry faces a critical shortage of experienced marine engineers. "
         "Consultancies that scale smart — with AI — will command premium rates."),
        ("Oceans5 is positioned at the centre",
         "Norway is the engineering talent hub for offshore energy. We combine that talent with AI efficiency."),
    ]
    for i, (t, b) in enumerate(market_cards):
        row, col = divmod(i, 2)
        x = Inches(7.5 + col * 2.9)
        y = Inches(1.8 + row * 2.5)
        add_rect(s, x, y, Inches(2.75), Inches(2.2), WHITE)
        add_tb(s, t, x + Inches(0.15), y + Inches(0.1), Inches(2.5), Inches(0.65),
               size=14, bold=True, color=TEAL, font_name="Montserrat")
        add_tb(s, b, x + Inches(0.15), y + Inches(0.75), Inches(2.5), Inches(1.3),
               size=10, color=TEXT_DARK, font_name="Inter")

    add_footer(s, "Oceans5 · Market Opportunity", dark=True)
    return s


def slide_10_vision(prs):
    s = add_blank(prs)
    add_rect(s, 0, 0, W, H, NAVY_DEEP)

    add_tag(s, "Vision 2028 — Nordstjernen", Inches(0.6), Inches(0.4))
    add_tb(s, "From resource supplier to premium partner",
           Inches(0.6), Inches(0.7), Inches(12), Inches(0.9),
           size=30, bold=True, color=WHITE, font_name="Montserrat")
    accent_bar(s, Inches(0.6), Inches(1.55))

    roadmap = [
        ("2026", "Foundation",
         "AI strategy adopted. QMS piloted. First work packages live. 90-day pilot rhythm established."),
        ("2027", "Scale",
         "Winning pilots scaled. First published reference cases. FEED and survey delivery live. "
         "ISO 9001 compliant. AI measurably embedded. Customer base doubled."),
        ("2028", "Leadership",
         "30–40% revenue from packages & services. 5%+ EBITDA improvement. ISO 9001 certified. "
         "Nordic market leader in AI-driven marine engineering."),
    ]
    for i, (year, title, body) in enumerate(roadmap):
        y = Inches(1.85 + i * 1.45)
        add_rect(s, Inches(0.4), y, Inches(0.02), Inches(1.25),
                 TEAL_BRIGHT if i == 2 else TEAL if i == 1 else GREY_MID)
        add_tb(s, year, Inches(0.55), y, Inches(1.1), Inches(0.6),
               size=20, bold=True, color=TEAL_BRIGHT, font_name="Montserrat")
        add_card(s, Inches(1.7), y, Inches(5.6), Inches(1.25))
        add_tb(s, title, Inches(1.85), y + Inches(0.1), Inches(5.35), Inches(0.45),
               size=13, bold=True, color=WHITE, font_name="Montserrat")
        add_tb(s, body, Inches(1.85), y + Inches(0.5), Inches(5.35), Inches(0.7),
               size=10, color=GREY_LIGHT, font_name="Inter")

    # Revenue mix chart
    chart_data = ChartData()
    chart_data.categories = ['Work packages & services', 'Resource supply']
    chart_data.add_series('2028 Target', (35, 65))
    chart = s.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(7.4), Inches(1.8), Inches(5.4), Inches(3.5),
        chart_data
    )
    try:
        pt = chart.chart.plots[0]
        pt.series[0].points[0].format.fill.solid()
        pt.series[0].points[0].format.fill.fore_color.rgb = TEAL_BRIGHT
        pt.series[0].points[1].format.fill.solid()
        pt.series[0].points[1].format.fill.fore_color.rgb = GREY_MID
    except Exception:
        pass

    add_rect(s, Inches(7.2), Inches(5.45), Inches(5.7), Inches(1.0), RGBColor(0x00, 0x35, 0x50))
    add_tb(s, '"When the project has a hard part — the client calls us first."',
           Inches(7.35), Inches(5.5), Inches(5.5), Inches(0.6),
           size=13, bold=True, color=WHITE, font_name="Montserrat")
    add_tb(s, "Oceans5 Vision Statement", Inches(7.35), Inches(6.05), Inches(5.5), Inches(0.35),
           size=9, color=GREY_LIGHT, font_name="Inter")

    add_footer(s, "Oceans5 · Vision 2028")
    return s


def slide_11_why(prs):
    s = add_blank(prs)
    add_rect(s, 0, 0, W, H, NAVY_DEEP)

    add_tag(s, "Why Oceans5", Inches(0.6), Inches(0.4))
    add_tb(s, "The combination no one else has yet assembled",
           Inches(0.6), Inches(0.7), Inches(12), Inches(0.9),
           size=30, bold=True, color=WHITE, font_name="Montserrat")
    accent_bar(s, Inches(0.6), Inches(1.55))

    # Table header
    headers = ["Capability", "Large Consultancy", "Traditional Marine Firm", "Oceans5"]
    col_x = [Inches(0.4), Inches(3.5), Inches(6.4), Inches(9.7)]
    col_w = [Inches(3.0), Inches(2.8), Inches(3.2), Inches(3.4)]

    for i, (h, cx) in enumerate(zip(headers, col_x)):
        clr = TEAL_BRIGHT if i == 3 else GREY_LIGHT
        add_tb(s, h, cx, Inches(1.75), col_w[i], Inches(0.4),
               size=11, bold=True, color=clr, font_name="Montserrat")

    add_rect(s, Inches(0.4), Inches(2.12), Inches(12.5), Pt(1.5), TEAL)

    rows = [
        ("Senior hands-on expertise",       "Junior delivery",    "✓ Yes",       "✓ Senior, always"),
        ("Full-scope: engineering + ops",    "Engineering only",   "Ops only",    "✓ Both"),
        ("AI embedded in delivery",          "Pilots only",        "Not yet",     "✓ Daily, measured"),
        ("Vendor-independent advice",        "Conflict of interest","✓ Usually",  "✓ Always"),
        ("Work packages (outcome-based)",    "✓ Yes",              "Hourly only", "✓ Growing"),
        ("Nordic offshore expertise",        "Generic",            "✓ Yes",       "✓ Core focus"),
        ("Agility & responsiveness",         "Slow, hierarchical", "Moderate",    "✓ Built-in"),
    ]
    for r, row in enumerate(rows):
        y = Inches(2.25 + r * 0.6)
        if r % 2 == 0:
            add_rect(s, Inches(0.4), y, Inches(12.5), Inches(0.58), RGBColor(0x10, 0x23, 0x38))
        for i, (cell, cx) in enumerate(zip(row, col_x)):
            is_check = cell.startswith("✓")
            clr = TEAL_BRIGHT if (i == 3 and is_check) else WHITE if is_check else GREY_MID
            add_tb(s, cell, cx, y + Inches(0.08), col_w[i], Inches(0.45),
                   size=10, color=clr, font_name="Inter")

    add_footer(s, "Oceans5 · Why Choose Us")
    return s


def slide_12_cta(prs):
    s = add_blank(prs)
    add_rect(s, 0, 0, W, H, NAVY_DEEP)
    add_rect(s, Inches(7), 0, Inches(6.333), H, RGBColor(0x00, 0x5A, 0x75))

    add_tb(s, "Oceans5", Inches(3.0), Inches(0.8), Inches(7.3), Inches(1.2),
           size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Montserrat")
    add_rect(s, Inches(5.1), Inches(1.9), Inches(3.2), Pt(3), TEAL_BRIGHT)

    add_tb(s, "Ready to deliver your most complex\nproject — better.",
           Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.5),
           size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Montserrat")

    add_tb(s, ("Let's start with a conversation about your project challenges.\n"
               "No generic slide deck — a real technical discussion about how we can help."),
           Inches(2.5), Inches(3.8), Inches(8.3), Inches(1.0),
           size=13, color=GREY_LIGHT, align=PP_ALIGN.CENTER, font_name="Inter")

    contacts = [
        ("📍 Address", "Drammensveien 126A\n0277 Oslo, Norway"),
        ("✉️ Email", "post@oceans5.no\noceans5.no"),
        ("📞 CEO Direct", "Peter-Emil Johannessen\n+47 412 69 164"),
    ]
    for i, (t, b) in enumerate(contacts):
        x = Inches(1.3 + i * 3.7)
        add_card(s, x, Inches(5.0), Inches(3.4), Inches(1.9))
        add_tb(s, t, x + Inches(0.15), Inches(5.1), Inches(3.1), Inches(0.45),
               size=11, bold=True, color=TEAL_BRIGHT, align=PP_ALIGN.CENTER, font_name="Montserrat")
        add_tb(s, b, x + Inches(0.15), Inches(5.55), Inches(3.1), Inches(1.15),
               size=11, color=GREY_LIGHT, align=PP_ALIGN.CENTER, font_name="Inter")

    add_footer(s, "© 2026 Oceans5 Group AS · All rights reserved", "oceans5.no")
    return s


# ═══════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════
def main():
    prs = new_prs()
    slide_01_cover(prs)
    slide_02_challenge(prs)
    slide_03_who_we_are(prs)
    slide_04_services(prs)
    slide_05_ai_advantage(prs)
    slide_06_track_record(prs)
    slide_07_approach(prs)
    slide_08_team(prs)
    slide_09_market(prs)
    slide_10_vision(prs)
    slide_11_why(prs)
    slide_12_cta(prs)

    prs.save(OUT)
    print(f"Saved: {OUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
